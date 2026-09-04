from pptx import Presentation

# Inicializa a apresentação
prs = Presentation()

# Função auxiliar para criar slides padronizados
def criar_slide(layout_idx, titulo, topicos):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = titulo
    if topicos and len(slide.placeholders) > 1:
        corpo = slide.placeholders[1].text_frame
        corpo.text = topicos[0]
        for topico in topicos[1:]:
            p = corpo.add_paragraph()
            p.text = topico
            p.level = 0
    return slide

# ==========================================
# SLIDE 1: Capa
# ==========================================
slide_capa = prs.slides.add_slide(prs.slide_layouts[0])
slide_capa.shapes.title.text = "CloudPulse V20"
slide_capa.placeholders[1].text = "Inteligência Competitiva com IA Multiagente\n\nFabrício Pinheiro Santos\nUniversidade Anhembi Morumbi"

# ==========================================
# SLIDE 2: O Problema
# ==========================================
criar_slide(1, "O Custo do Ponto Cego Competitivo", [
    "❌ Monitoramento manual: Horas perdidas rastreando concorrentes.",
    "❌ Dados espalhados: Informações perdidas em planilhas e e-mails.",
    "❌ Lentidão comercial: Reação demorada a promoções rivais.",
    "❌ Retrabalho constante: Zero automação no processo analítico."
])

# ==========================================
# SLIDE 3: Nossa Solução
# ==========================================
criar_slide(1, "CloudPulse V20: O Antídoto", [
    "🤖 IA Multiagente: Especialistas virtuais delegando tarefas.",
    "🧠 RAG Inteligente: Respostas precisas baseadas em dados reais.",
    "📊 Dashboard Web: O mercado inteiro em uma única tela.",
    "⚙️ Coleta Automática: Crawler invisível e implacável.",
    "📈 Relatórios Executivos: Insights diretos para tomada de decisão."
])

# ==========================================
# SLIDE 4: Benefícios e Impacto
# ==========================================
criar_slide(1, "Valor Real para o Negócio", [
    "✅ Atualização Automática: Dados do mercado sempre frescos.",
    "✅ Alertas Inteligentes: Notificação imediata de mudanças.",
    "✅ Dashboard em Tempo Real: Agilidade na ponta da linha.",
    "✅ Busca Vetorial: Contexto exato em milissegundos.",
    "✅ Escalabilidade Serverless: Cresce com a demanda, sem travar.",
    "✅ Decisão sem Alucinação: Dados verificados pela arquitetura."
])

# ==========================================
# SLIDE 5: Arquitetura
# ==========================================
criar_slide(1, "Arquitetura Cloud & Serverless", [
    "O motor por trás da plataforma:",
    "• Frontend e Gatilho: CloudFront ➔ S3 ➔ API Gateway",
    "• Orquestração e Processamento: AWS Lambda ➔ Step Functions",
    "• Cérebro da Inteligência: Amazon Bedrock ➔ OpenSearch",
    "• Persistência e Notificação: DynamoDB ➔ SES",
    "• Interface do Cliente: Dashboard Web em Tempo Real"
])

# ==========================================
# SLIDE 6: Fluxograma
# ==========================================
criar_slide(1, "O Fluxo Multiagente", [
    "1. Usuário faz uma pergunta no Dashboard.",
    "2. API Gateway roteia para o Agente Supervisor.",
    "3. O Supervisor delega para os Especialistas (Crawler, Analista).",
    "4. Sistema executa RAG (Retrieval-Augmented Generation).",
    "5. Resposta executiva e relatório são entregues na tela."
])

# ==========================================
# SLIDE 7: Demonstração Ao Vivo
# ==========================================
criar_slide(1, "Demonstração Ao Vivo", [
    "O motor rodando na prática:",
    "► Login e Acesso Seguro",
    "► Dashboard de Monitoramento",
    "► Pesquisa Interativa",
    "► Geração de Relatório via IA",
    "► O Fim das Apresentações Estáticas."
])

# ==========================================
# SLIDE 8: Nossos Diferenciais
# ==========================================
criar_slide(1, "Por que o CloudPulse Domina?", [
    "✔ IA Multiagente vs Chatbots Comuns",
    "✔ Agente Supervisor de Contexto",
    "✔ RAG Híbrido de altíssima precisão",
    "✔ Totalmente Event-Driven e Serverless",
    "✔ Segurança com Autenticação (Cognito)",
    "✔ Arquitetura pronta para SaaS (Software as a Service)"
])

# ==========================================
# SLIDE 9: Evoluções Futuras
# ==========================================
criar_slide(1, "Roadmap de Expansão (SaaS)", [
    "• Self-Correction Agent (Validação autônoma dos dados)",
    "• Chat com Documentos (Análise de contratos e PDFs)",
    "• IA Multimodal (Avaliação visual de banners concorrentes)",
    "• Integração B2B (Alertas via WhatsApp, Slack e CRM)",
    "• Multi-Tenant (Painel administrativo e analytics global)"
])

# ==========================================
# SLIDE 10: Contato
# ==========================================
criar_slide(1, "Obrigado!", [
    "A inteligência do mercado agora está nas suas mãos.",
    "",
    "[ INSERIR QR CODE DO SISTEMA AQUI ]",
    "",
    "GitHub: 1205fa/CloudPulse-Dashboard-V20",
    "LinkedIn: in/fabricio-pinheiro-santos",
    "Contato: fcalberico77@gmail.com"
])

# Salva o arquivo final
prs.save("CloudPulse_Apresentacao.pptx")
print("✅ Apresentação gerada com sucesso: CloudPulse_Apresentacao.pptx")
